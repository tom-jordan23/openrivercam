#!/usr/bin/env python3
"""
build_deck.py — Companion slide deck for the replication recommendations report.

Derived from docs/REPLICATION_RECOMMENDATIONS.md. The report is the source of
record; this script holds only the slide-level summary of it. Scope matches the
report: technology only. How the output is applied is for IPB, BHLK and their
partners, and nothing here draws a hydrology conclusion.

Usage:
    ./build_deck.py                                  # neutral template
    ./build_deck.py --template /path/to/amcross.pptx # branded
    ./build_deck.py -o pdf/REPLICATION_BRIEFING.pptx

Requires python-pptx:
    uv pip install --python .venv-pdf/bin/python python-pptx
"""

import argparse
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
PHOTOS = HERE.parent / "build_photos"
IMAGES = HERE / "images"

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import PP_PLACEHOLDER
except ImportError:
    sys.exit("python-pptx not installed. See the header of this script.")

DOC_TITLE = "OpenRiverCam in Indonesia"
SUBTITLE = "What the pilot taught us, and the path forward\nResponse to the PMI / IPB / BHLK meeting, Sukabumi, 21 August 2026"
# Movement convention: all parties equal. Same order, same weight, everywhere.
PARTNERS = ("Palang Merah Indonesia  ·  Institut Pertanian Bogor  ·  "
            "Balai Hidrologi dan Lingkungan Keairan  ·  American Red Cross")
LOGOS = [
    ("pmi.png", "Palang Merah Indonesia"),
    ("ipb.png", "Institut Pertanian Bogor"),
    ("bhlk.png", "Balai Hidrologi dan Lingkungan Keairan"),
    ("amcross.png", "American Red Cross"),
]
FOOTER = "ORC Indonesia — PMI · IPB · BHLK · American Red Cross"
STATUS = "Draft for internal review — not yet circulated"

INK = RGBColor(0x1A, 0x1A, 0x1A)
SECOND = RGBColor(0x42, 0x41, 0x3E)
MUTED = RGBColor(0x6B, 0x6A, 0x66)

# The Red Cross templates ship with example slides. Anything shipped in the
# repo alongside this script names the template it was built from.
DEFAULT_TEMPLATE = ("/home/tjordan/code/templates/AmCross/English PowerPoint "
                    "Templates/502201-04 Red Cross Classic Template FINAL.pptx")


# ── Slide content ────────────────────────────────────────────────
# Each entry is (kind, payload). Kinds: title, section, bullets, table.
# Bullet tuples are (level, text). Level 0 is a top-level point.

DECK = [
    ("title", {}),

    ("bullets", {
        "title": "Thank you",
        "bullets": [
            (0, "None of this exists without the people who built it and looked "
                "after it."),
            (0, "PMI volunteers and staff at Sukabumi and Jakarta gave their time "
                "to a system that was new to them and that did not always work."),
            (0, "IPB re-surveyed the site with a total station after our own "
                "survey failed twice. That survey is what the station runs on "
                "today."),
            (0, "BHLK brought the standards knowledge, and the offer of server "
                "capacity, that make a pilot possible at all."),
            (0, "What follows is largely a list of things we would do "
                "differently. Every one was learned because someone did the work "
                "that made it visible."),
        ],
    }),

    ("bullets", {
        "title": "The pilot, and what we think comes next",
        "bullets": [
            (0, "The pilot has been a good experience. A station in the water, "
                "four months of behaviour to learn from, and a working "
                "relationship between PMI, IPB and BHLK that did not exist "
                "before."),
            (0, "Our recommendation is that the next step is yours."),
            (1, "Not because the partnership has run its course, but because of "
                "what the pilot produced: knowledge, not a body of "
                "measurements."),
            (1, "Knowledge is worth more to you as an input to your own approach "
                "than as a set of parts to copy."),
            (0, "We would rather help you start your own design than hand you "
                "ours to maintain."),
        ],
    }),

    ("bullets", {
        "title": "What this briefing covers, and what it leaves to you",
        "bullets": [
            (0, "The technology only: the station as built, what four months in "
                "the field revealed, and what we would change before more units "
                "are constructed."),
            (0, "Not in scope: what the data is fit to support, or what accuracy "
                "any application requires."),
            (1, "Those judgements belong with IPB, BHLK and their federal "
                "partners."),
            (1, "Measurement requirements shown here are recorded from you, not "
                "proposed by us."),
            (0, "Written to give an accurate basis for your decision rather than "
                "an encouraging one."),
        ],
    }),

    ("bullets", {
        "title": "Three offers were made on 21 August. We welcome all three.",
        "bullets": [
            (0, "Build one to three stations as a pilot."),
            (1, "Supported. What follows is what we would change first."),
            (0, "Provide server capacity for the data."),
            (1, "Well matched to the pilot, and holding the data in Indonesia "
                "carries a clear benefit."),
            (0, "Move the present site to open ground."),
            (1, "Supported, and the field record adds independent evidence for "
                "it."),
        ],
    }),

    ("bullets", {
        "title": "Why a low-cost station is worth studying",
        "bullets": [
            (0, "A river gauge is only useful if there are enough of them. How "
                "many an agency can afford decides how much of a catchment it "
                "can observe."),
            (0, "This station: USD 1,340 in materials — electronics and "
                "enclosure only. Sukabumi already had its panel and battery, so a "
                "new solar site would add an array. A mains site would not."),
            (0, "Lowest-cost automatic water-level station on the government "
                "e-catalogue: about USD 3,600 before VAT, measuring stage only."),
            (1, "Not like for like — that is a supported commercial instrument, "
                "this is a pilot-stage assembly."),
            (0, "At this price a basin authority can consider a network where it "
                "might otherwise consider a single gauge. That possibility is "
                "the whole of the argument."),
        ],
    }),

    ("figure", {
        "title": "What is being offered for study",
        "image": "fig1_system.png",
        "alt": "Flow diagram. A camera on a pole at the river sends video to a "
               "station computer, which works out a water level and a discharge "
               "figure; both travel over the mobile network to a server that "
               "keeps the record. Beneath the computer are the three things it "
               "depends on: a solar panel and battery, a scheduler that wakes it "
               "every 30 minutes, and a rain gauge that keeps recording while the "
               "station sleeps.",
    }),

    ("photos", {
        "title": "Built from commodity parts, so it can be repaired locally",
        "photos": [
            (PHOTOS / "sukabumi" / "IMG_0048.png",
             "The parts for one station, before assembly.",
             "Components laid out on a workbench: enclosure mounting plate, two "
             "lengths of DIN rail, three fuse holders, the camera, the rain gauge "
             "dome, a small computer with its screw-terminal riser, a power "
             "converter, terminal blocks, a relay board, the mobile data modem "
             "and its antenna."),
            (IMAGES / "sukabumi" / "complete-system-before-power.png",
             "The same parts wired onto the plate, under bench power.",
             "The same components assembled onto the mounting plate and wired, "
             "with a bench power supply alongside reading 12.08 volts."),
        ],
        "note": "No soldering, no fabrication, nothing single-source.",
    }),

    ("bullets", {
        "title": "The evidence is thin, and we ask you to treat it as thin",
        "bullets": [
            (0, "One station, on one river, watched across part of one season — "
                "16 April to 28 August 2026."),
            (0, "A second station was built but never installed."),
            (0, "What the deployment has produced so far is knowledge about the "
                "design, not a body of measurements."),
            (0, "It is a volunteer pilot and should be read as one."),
            (1, "Never built to production standards of availability or record "
                "continuity. Measuring it against an industrial instrument would "
                "be the wrong test, and unfair to the people who kept it "
                "running."),
            (1, "The useful question is narrower: what does this design make "
                "hard, and what should be built differently."),
        ],
    }),

    ("bullets", {
        "title": "Three design gaps, and what we would build instead",
        "bullets": [
            (0, "The station cannot report its own condition. When it stopped, it "
                "stopped quietly."),
            (1, "Instead: report state on every waking, and make any mode that "
                "suppresses data expire on its own and raise an alert. R4, R5."),
            (0, "Nothing reconciles what was recorded against what arrived."),
            (1, "Instead: compare the station's record against the server's on a "
                "schedule. Small piece of software; should have been in version "
                "one. R7."),
            (0, "Water level is read from the image, and daylight defeats it."),
            (1, "Instead: an independent water-level reference, so the "
                "measurement does not rest on one optical method succeeding. "
                "R1."),
        ],
    }),

    ("bullets", {
        "title": "A twenty-five dollar decision caused most of the outages",
        "bullets": [
            (0, "No single step in the chain looks like a mistake."),
            (1, "The USB storage drive caused a driver fault at boot, so it was "
                "removed rather than fixed."),
            (1, "That left the SD card as the only volume — small enough to sit "
                "at the level where it deletes old recordings."),
            (1, "That caused processing to fail on 43% of videos."),
            (1, "The station shuts down after processing finishes, so a failed "
                "run never shut down: awake to a 25-minute backstop instead of "
                "two minutes, about twelve times the energy."),
            (1, "Repeated across a night, that flattened the battery. A missed "
                "wake then left the next-startup alarm in the past."),
            (0, "Two lessons: shutdown and startup are split across two systems "
                "and neither owns the cycle; and cheap parts can carry expensive "
                "operating costs."),
            (0, "A common root: the computer's own clock battery connector broke "
                "on both boards. Losing it is why a separate scheduling board was "
                "fitted, which is what split the cycle. R37"),
        ],
    }),

    ("figure", {
        "title": "The station measures water level reliably only at night",
        "image": "fig2_optical.png",
        "alt": "Two panels from 200 captures. The upper panel counts captures by "
               "hour of day, split into those that produced a water level and "
               "those rejected; every rejection falls between 06:00 and 19:00, "
               "with peaks mid-morning and mid-afternoon and fewer in the early "
               "afternoon. The lower panel shows how confident each reading was "
               "against the threshold at which a water level is accepted.",
        "note": "The two peaks and the early-afternoon dip are the pattern a "
                "sun-angle effect produces.",
    }),

    ("groups", {
        "title": "What we recommend",
        "note": "Full text in the working list; the report carries an index of "
                "all of them.",
    }),

    ("bullets", {
        "title": "The changes with the greatest effect",
        "bullets": [
            (0, "Use an industrial Pi carrier with NVMe, a protected clock and an "
                "integrated UPS. Most of the failures above meet at one place, "
                "and this removes them in a single part. R37"),
            (0, "Build the test station first, so a fault can be reproduced and a "
                "fix tried before it goes to a river. R36"),
            (0, "Build monitoring for a fleet, not for one station. This is the "
                "highest-value work on the list. R4–R7"),
            (0, "Fit an independent water-level reference, and plan the survey as "
                "skilled work from the start. R1, R2, R19, R20"),
            (0, "Give one process control of the whole sleep and wake cycle, and "
                "make shutdown happen on a timer. R10, R12"),
            (0, "Make mains power the default; use solar only where mains is not "
                "available. R11"),
            (0, "Choose the site before anything else, and confirm permission in "
                "writing before building for it. R16–R18"),
        ],
    }),

    ("figure", {
        "title": "Camera at the river, computer indoors (R8)",
        "image": "fig3_configurations.png",
        "alt": "Side-by-side comparison. On the left, the arrangement as built: "
               "camera, computer, modem and power system all in an enclosure at "
               "the riverbank, so all of it is in the weather and service means "
               "opening the enclosure on site. On the right, the arrangement "
               "proposed for the pilot units: only the camera stays at the river "
               "and the computer runs indoors at a BHLK or IPB facility over a "
               "network link.",
        "note": "Designed and not yet field tested — the pilot's first "
                "experiment, not a proven alternative.",
    }),

    ("bullets_photo", {
        "title": "One lesson about method",
        "bullets": [
            (0, "The budget was applied part by part: the cheapest item meeting "
                "each requirement."),
            (0, "It prices each part against its specification sheet, not against "
                "what that part's limits cost the rest of the system."),
            (0, "The camera meets every line of its specification, and still "
                "costs a fifth of the video quality and 30–60 s of battery on "
                "every waking."),
            (0, "Its light fires whenever it starts and cannot be turned off — 48 "
                "flashes a day at the present site."),
        ],
        "photo": IMAGES / "components" / "annke_c1200_camera.png",
        "caption": "Capable hardware, running a reseller's version of the "
                   "manufacturer's software.",
        "alt": "The camera as delivered, on a workbench mat with its mounting "
               "hardware, waterproof cable boot and printed manual.",
    }),

    ("bullets", {
        "title": "Build the test station first",
        "bullets": [
            (0, "Everything we know about the failures was diagnosed on a solar "
                "station, on a river, awake for tens of seconds, that we could "
                "not touch."),
            (1, "We could not reproduce a fault, test a fix before committing it "
                "remotely, or tell a working change from a quiet week. Several "
                "diagnoses took months for that reason alone."),
            (0, "The Jakarta station was meant to be that test station and never "
                "became one. Its absence cost more than the site it was built "
                "for."),
            (0, "A test station is not a spare: mains-powered, always on, "
                "somewhere someone can watch it, open it and break it "
                "deliberately."),
            (0, "We would treat it as the first station a pilot builds, not the "
                "last."),
        ],
    }),

    ("bullets", {
        "title": "What is worth keeping",
        "bullets": [
            (0, "The five constraints: commodity parts, no soldering, no "
                "specialist skills, common tools, five-minute replacement."),
            (0, "The spare switched outputs. The relay module has four channels; "
                "one drives the camera and three are deliberately left free, "
                "wired and ready."),
            (1, "That is what lets a station drive a siren, a beacon or an "
                "alerting relay without reopening the design — and what turns a "
                "measurement station into something a community can act on."),
            (1, "We would ask that any replication keeps that spare capacity "
                "rather than removing it to save a few dollars."),
            (0, "The factory-sealed camera, which removed the humidity failure "
                "that killed the previous unit."),
            (0, "The recovery kit, documentation in both languages, and spares at "
                "the local PMI chapter."),
        ],
    }),

    ("bullets", {
        "title": "Building and looking after the units in Indonesia",
        "bullets": [
            (0, "The design was made to be built by people who are not "
                "electronics specialists, with tools they already own."),
            (0, "That intent only becomes real if the parts can be bought "
                "locally and the documentation supports substitution — R3."),
            (0, "Spares held at the local PMI chapter turn a failure into a part "
                "swap rather than a shipment."),
            (0, "Where the field unit is a standard security camera, support "
                "falls within a supply chain that already exists across "
                "Indonesia."),
            (0, "We would rather help build the capacity to maintain these "
                "stations than remain the place they are sent when they break."),
        ],
    }),

    ("bullets", {
        "title": "Hosting the data in Indonesia",
        "bullets": [
            (0, "BHLK's offer suits the pilot, and holding the data in Indonesia "
                "carries a clear benefit for a government-partnered deployment."),
            (0, "Two constraints to design for rather than discover:"),
            (1, "The video store and the database must share one filesystem. Plan "
                "storage as a single volume."),
            (1, "Server and station software versions are tied together, and a "
                "remote station cannot be upgraded on demand."),
            (0, "We would suggest running both instances in parallel until the "
                "new one has completed a full operating cycle, including an "
                "upgrade."),
            (0, "Worth agreeing in writing: where the authoritative copy sits, "
                "who administers it, who has access, and the retention policy."),
        ],
    }),

    ("bullets", {
        "title": "Choosing the sites",
        "bullets": [
            (0, "An open site helps with three problems at once: satellite "
                "positioning, the sun-water-camera angle, and the view across "
                "the section."),
            (0, "A move costs a fresh survey and a fresh calibration, not only a "
                "relocation. The survey is the expensive part."),
            (0, "Site permission should be settled before a unit is built for a "
                "particular site."),
            (1, "A complete station was built and tested for an intended "
                "Jakarta site. Permission fell through during the April visit."),
            (0, "That station is still available. We suggest it serves best as a "
                "study and test unit."),
            (1, "Open it, trace it, power it up, take it apart — and install it "
                "somewhere convenient and local if that is useful, so it can be "
                "exercised against real water."),
            (1, "Not as an operational station with expectations of availability "
                "and consistent data. It carries the design this briefing "
                "recommends changing."),
            (0, "We would not build to a site again before you tell us the "
                "permission is in place."),
        ],
    }),

    ("bullets", {
        "title": "Working together",
        "bullets": [
            (0, "BHLK — data processing, conformance with standards, and the "
                "route to acceptance within PUPR data systems. Offers server "
                "capacity."),
            (0, "IPB — design, calibration methodology, training material, and "
                "the development path for future sensor types."),
            (0, "PMI — user of the information, and the operational side: "
                "installation, maintenance, siting, response, and spares."),
            (0, "Offered for your consideration: write the split into the "
                "collaboration agreement so it survives staff changes, and keep "
                "a light joint forum for what does not sort into one role."),
            (0, "The commitments this implies for PMI have not been discussed "
                "with PMI National Headquarters."),
        ],
    }),

    ("bullets", {
        "title": "Questions for consideration",
        "bullets": [
            (0, "Who should hold the Jakarta station as a study and test unit? "
                "It was to transfer to IPB for deployment. Whether IPB or BHLK "
                "is better placed to hold it is for the two of you."),
            (0, "Can a staff gauge be read from the camera image accurately "
                "enough? If it can, R1 needs no separate sensor."),
            (0, "Two of the thirteen interruptions remain unexplained."),
            (0, "Does the recovery-voltage threshold bound how long an "
                "interruption lasts, or hold the station off?"),
            (0, "How does the velocity measurement perform across the surface "
                "conditions the pilot sites present?"),
            (0, "Absolute discharge accuracy at Sukabumi, unresolved pending the "
                "survey."),
            (0, "We expect your work to change some of what is written here. "
                "That is the point of handing it over rather than handing it "
                "down."),
        ],
    }),

    ("bullets", {
        "title": "What we can offer from here",
        "bullets": [
            (0, "We are not proposing to build your stations."),
            (0, "What we can offer is the record: this briefing and its report, "
                "the appendix, the operator and assembly documentation in English "
                "and Bahasa Indonesia, the software, and the built station at "
                "Wisma PMI to take apart."),
            (0, "Beyond that, whatever is useful — reviewing a design, answering "
                "a question about something that surprised us, looking at data "
                "that does not behave."),
            (0, "Sukabumi will keep running, and we will keep reporting what it "
                "does, including the parts that go wrong."),
            (0, "The pilot found the problems while they were still cheap, and it "
                "introduced three organisations to one another. The next station "
                "in this story should be one you designed."),
        ],
    }),
]


# ── Rendering ────────────────────────────────────────────────────
# Written against the American Red Cross Classic template, whose slides are
# 10 x 5.62 in with a 3.35 in body. Everything is positioned from the layout's
# own placeholders rather than from fixed offsets, so a different template's
# geometry is followed rather than overridden.

CHROME = (PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER, PP_PLACEHOLDER.DATE)


def find_layout(prs, *names, fallback=1):
    """Look a layout up by name across the template's masters, then by index."""
    wanted = [n.lower() for n in names]
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name.lower() in wanted:
                return layout
    layouts = prs.slide_layouts
    return layouts[fallback] if fallback < len(layouts) else layouts[0]


def is_chrome(shape):
    """Footer, date and slide-number placeholders, which the template manages."""
    if not shape.is_placeholder:
        return False
    try:
        return shape.placeholder_format.type in CHROME
    except ValueError:
        return False


def body_placeholders(slide):
    """Content placeholders, in layout order, excluding title and chrome.

    python-pptx hands back a fresh proxy on each access, so the title has to be
    identified by its underlying element rather than by object identity.
    """
    title = slide.shapes.title
    title_el = title._element if title is not None else None
    out = []
    for ph in slide.placeholders:
        if ph._element is title_el or is_chrome(ph) or not ph.has_text_frame:
            continue
        out.append(ph)
    return out


def layout_chrome(slide):
    """Footer and slide-number rectangles as the template's layout defines them.

    python-pptx does not clone footer, date or slide-number placeholders onto a
    new slide, so their geometry has to be read back off the layout in order to
    put anything in those positions.
    """
    footer = number = None
    for ph in slide.slide_layout.placeholders:
        try:
            kind = ph.placeholder_format.type
        except ValueError:
            continue
        rect = (ph.left, ph.top, ph.width, ph.height)
        if None in rect:
            continue
        if kind == PP_PLACEHOLDER.FOOTER:
            footer = rect
        elif kind == PP_PLACEHOLDER.SLIDE_NUMBER:
            number = rect
    return footer, number


def clear_empty_placeholders(slide):
    """Drop placeholders left untouched, so template prompt text is not printed.

    Chrome is left alone — removing the slide-number placeholder would take the
    template's numbering with it.
    """
    for shape in list(slide.placeholders):
        if is_chrome(shape) or not shape.has_text_frame:
            continue
        if shape.text_frame.text.strip() == "":
            shape._element.getparent().remove(shape._element)


def write(tf, lines, wrap=True):
    """Fill a text frame with (text, size, colour, level) tuples."""
    tf.word_wrap = wrap
    for i, (text, size, color, level) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color


def style_title(slide, text, size=None):
    """Set the slide title, sized so it does not run into the body.

    Template title placeholders are laid out for a short line; these titles are
    sentences, so the size is chosen from the length.
    """
    title = slide.shapes.title
    if title is None:
        return None
    if size is None:
        size = 20 if len(text) <= 55 else (18 if len(text) <= 80 else 16)
    write(title.text_frame, [(text, size, INK, 0)])
    return title


def content_box(slide, prs):
    """Left, top, width and available height for a shape under the title."""
    title = slide.shapes.title
    left = Inches(0.5)
    top = Inches(1.31)
    width = prs.slide_width - Inches(1.0)
    if title is not None and title.top is not None:
        left = title.left
        width = title.width
        top = title.top + title.height + Inches(0.1)

    # Stop clear of the footer band, which the template puts at 5.10 in.
    bottom = prs.slide_height - Inches(0.62)
    return left, top, width, max(bottom - top, Inches(1.0))


def add_title_slide(prs):
    lines = [
        (DOC_TITLE, 20, INK, 0),
        (SUBTITLE, 12, MUTED, 0),
        (PARTNERS, 10, MUTED, 0),
        ("2026-08-31  ·  " + STATUS, 9, MUTED, 0),
    ]

    if not logo_files():
        # No partner artwork: fall back to the template's own title slide.
        layout = find_layout(prs, "title slide", "title slide 1", "title",
                             fallback=0)
        slide = prs.slides.add_slide(layout)
        title = style_title(slide, DOC_TITLE, size=20)
        bodies = body_placeholders(slide)
        if title is None and bodies:
            write(bodies[0].text_frame, lines)
        elif bodies:
            write(bodies[0].text_frame, lines[1:])
        clear_empty_placeholders(slide)
        return slide

    # With artwork, the template's title slide is not usable: it centres the
    # American Red Cross mark alone above the title. Title Only is a blank
    # canvas, so the slide is composed here — logo row first, then the title
    # block beneath it, with all four organisations at one size in one row.
    layout = find_layout(prs, "title only", fallback=12)
    slide = prs.slides.add_slide(layout)
    cover_master_logo(slide, prs)

    place_logo_row(slide, prs, Inches(0.8), Inches(0.72),
                   prs.slide_width - Inches(1.6), Inches(0.62))

    title = slide.shapes.title
    if title is not None:
        title.left, title.top = Inches(0.5), Inches(1.95)
        title.width, title.height = prs.slide_width - Inches(1.0), Inches(0.85)
        write(title.text_frame, [(DOC_TITLE, 22, INK, 0)])

    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.95),
                                   prs.slide_width - Inches(1.8), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    write(tf, lines[1:])
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER

    clear_empty_placeholders(slide)
    return slide


def add_section_slide(prs, title, kicker=""):
    layout = find_layout(prs, "subsection", "section header", "divider slide: text only",
                         "title only", fallback=2)
    slide = prs.slides.add_slide(layout)
    style_title(slide, title, size=20)

    if kicker:
        bodies = body_placeholders(slide)
        if bodies:
            write(bodies[0].text_frame, [(kicker, 11, MUTED, 0)])
        else:
            left, top, width, _ = content_box(slide, prs)
            box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
            write(box.text_frame, [(kicker, 11, MUTED, 0)])

    clear_empty_placeholders(slide)
    return slide


def add_bullets_slide(prs, title, bullets):
    layout = find_layout(prs, "title and content", "content", "title and body", fallback=1)
    slide = prs.slides.add_slide(layout)
    style_title(slide, title)

    bodies = body_placeholders(slide)
    if bodies:
        body = bodies[0]
    else:
        left, top, width, height = content_box(slide, prs)
        body = slide.shapes.add_textbox(left, top, width, height)

    # The body is 3.35 in on this template. Size from the volume of text so a
    # long slide shrinks rather than overflowing the placeholder.
    weight = sum(1 + len(text) // 88 for _, text in bullets)
    top_size = 15 if weight <= 9 else (14 if weight <= 12 else 12.5)

    write(body.text_frame,
          [(text, top_size if level == 0 else top_size - 1.5,
            INK if level == 0 else MUTED, level)
           for level, text in bullets])

    clear_empty_placeholders(slide)
    return slide


def add_table_slide(prs, title, columns, rows, notes=None):
    layout = find_layout(prs, "title only", "title and content", fallback=5)
    slide = prs.slides.add_slide(layout)
    style_title(slide, title)
    left, top, width, avail = content_box(slide, prs)
    clear_empty_placeholders(slide)

    note_height = Inches(0.8) if notes else 0
    # Integer EMU throughout: python-pptx rejects floats, and the division below
    # produces one as soon as the row count does not divide evenly.
    row_h = int(min(Inches(0.34), (avail - note_height) / (len(rows) + 1)))
    height = row_h * (len(rows) + 1)

    shape = slide.shapes.add_table(len(rows) + 1, len(columns), left, top, width, height)
    table = shape.table
    for r in range(len(rows) + 1):
        table.rows[r].height = row_h

    # A narrow first column where it carries only a short key.
    if columns and columns[0] == "#":
        key = Inches(0.5)
        table.columns[0].width = key
        rest = (width - key) // (len(columns) - 1)
        for i in range(1, len(columns)):
            table.columns[i].width = rest

    def fill(cell, text, size, bold):
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        cell.text = text
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = INK

    for c, name in enumerate(columns):
        fill(table.cell(0, c), name, 10, True)
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            fill(table.cell(r, c), value, 9, False)

    if notes:
        box = slide.shapes.add_textbox(left, top + height + Inches(0.15),
                                       width, note_height)
        write(box.text_frame, [(n, 10, MUTED, 0) for n in notes])
    return slide


def recommendation_groups():
    """Group names and their recommendation numbers, read from the source list.

    The deck used to restate the recommendations, which meant editing them in two
    places and letting them drift. RECOMMENDATIONS.md is the source; this reads
    it. A bullet's bold statement can wrap across lines, so only the opening
    marker is matched.
    """
    import re
    src = HERE / "RECOMMENDATIONS.md"
    skip = {"Acknowledgement — for the report, not a recommendation",
            "Notes for you"}
    groups, current = [], None
    for line in src.read_text(encoding="utf-8").split("\n"):
        h = re.match(r"^## (.+)$", line)
        if h:
            name = h.group(1).strip()
            current = None if name in skip else name
            if current:
                groups.append((current, []))
            continue
        m = re.match(r"^- \*\*(R\d+)", line)
        if m and groups and current:
            groups[-1][1].append(m.group(1))
    return [(name, nums) for name, nums in groups if nums]


def logo_files():
    """The partner logos that are actually present, in fixed order.

    Missing artwork is not an error: the row lays out equally for whoever is
    there, and the organisations are carried by name regardless. See
    logos/README.md for the specification.
    """
    found = []
    for name, org in LOGOS:
        path = HERE / "logos" / name
        if path.is_file():
            found.append((path, org))
    return found


# The master's own mark, from the Classic template: 1.54 x 0.71 in at
# (0.34, 4.89). Covered rather than deleted — a shape inherited from the master
# cannot be removed from an individual slide.
MASTER_LOGO = (Inches(0.30), Inches(4.86), Inches(1.66), Inches(0.74))


def cover_master_logo(slide, prs):
    """Hide the template's single-organisation mark on a content slide."""
    from pptx.enum.shapes import MSO_SHAPE
    left, top, width, height = MASTER_LOGO
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.fill.background()

    # The theme gives every shape an outline and a drop shadow, which would draw
    # a visible box exactly where the point is that nothing should show. Those
    # come from the shape's <p:style> reference rather than from spPr, so the
    # whole style element goes; then an empty effect list makes the absence
    # explicit for renderers that would otherwise fall back to the theme.
    from pptx.oxml.ns import qn
    el = shape._element
    for style in el.findall(qn("p:style")):
        el.remove(style)
    spPr = el.spPr
    for tag in ("a:effectLst", "a:effectRef"):
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)
    spPr.append(spPr.makeelement(qn("a:effectLst"), {}))
    set_alt_text(shape, "")
    return shape


def place_logo_row(slide, prs, left, top, width, height):
    """Lay logos out normalised to a common height and evenly spaced.

    Equal billing is a layout property, not a courtesy: same height, same gaps,
    no primary position. Scaling by height rather than by area is what keeps a
    wide lockup and a square mark reading as equals.
    """
    from PIL import Image
    logos = logo_files()
    if not logos:
        return False

    widths = []
    for path, _ in logos:
        with Image.open(path) as im:
            iw, ih = im.size
        widths.append(int(iw * (height / ih)))

    gap = Inches(0.32)
    total = sum(widths) + gap * (len(logos) - 1)
    if total > width:                       # shrink to fit, keeping equal height
        scale = width / total
        height = int(height * scale)
        widths = [int(w * scale) for w in widths]
        gap = int(gap * scale)
        total = sum(widths) + gap * (len(logos) - 1)

    x = left + (width - total) // 2
    for (path, org), w in zip(logos, widths):
        pic = slide.shapes.add_picture(str(path), x, top, width=w, height=height)
        set_alt_text(pic, org + " logo")
        x += w + gap
    return True


def set_alt_text(shape, text):
    """Alternate text for a picture, read by screen readers in PowerPoint.

    python-pptx has no accessor for this, so it is written straight onto the
    non-visual drawing properties.
    """
    shape._element._nvXxPr.cNvPr.set("descr", text)


def place_image(slide, prs, path, box, alt):
    """Fit an image inside (left, top, width, height), preserving aspect and
    centring it in whichever direction has slack."""
    from PIL import Image
    left, top, width, height = box
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(width / iw, height / ih)
    w, h = int(iw * scale), int(ih * scale)
    pic = slide.shapes.add_picture(
        str(path), int(left + (width - w) / 2), int(top + (height - h) / 2),
        width=w, height=h)
    set_alt_text(pic, alt)
    return pic


def content_area(slide, prs):
    """The rectangle between the title and the template's footer band.

    The Classic template puts its footer rule at 5.10 in on a 5.62 in slide,
    so anything below about 4.95 in collides with it.
    """
    left, top, width, _ = content_box(slide, prs)
    bottom = prs.slide_height - Inches(0.80)
    return left, top, width, max(bottom - top, Inches(1.0))


def add_figure_slide(prs, title, image, alt, note=None):
    layout = find_layout(prs, "title only", "title and content", fallback=5)
    slide = prs.slides.add_slide(layout)
    t = style_title(slide, title)
    left, top, width, height = content_area(slide, prs)

    # The template sizes its title placeholder for two lines; these titles are
    # one. Reclaim the unused half by starting the image higher, rather than
    # resizing the placeholder — writing an explicit size onto an inheriting
    # placeholder also writes an offset of (0, 0) and moves the title.
    if t is not None and t.top is not None:
        gained = top - (t.top + Inches(0.78))
        if gained > 0:
            top -= gained
            height += gained

    clear_empty_placeholders(slide)

    note_h = Inches(0.46) if note else 0
    place_image(slide, prs, HERE / "figures" / "png" / image,
                (left, top, width, height - note_h), alt)
    if note:
        box = slide.shapes.add_textbox(left, top + height - note_h, width,
                                       note_h)
        box.text_frame.word_wrap = True
        write(box.text_frame, [(note, 10, MUTED, 0)])
    return slide


def add_photos_slide(prs, title, photos, note=None):
    """Two photographs side by side, each with its own short caption."""
    layout = find_layout(prs, "title only", "title and content", fallback=5)
    slide = prs.slides.add_slide(layout)
    style_title(slide, title)
    left, top, width, height = content_area(slide, prs)
    clear_empty_placeholders(slide)

    note_h = Inches(0.44) if note else 0
    cap_h = Inches(0.5)
    gap = Inches(0.3)
    cell_w = int((width - gap * (len(photos) - 1)) / len(photos))
    img_h = height - note_h - cap_h

    for i, (path, caption, alt) in enumerate(photos):
        x = left + i * (cell_w + gap)
        place_image(slide, prs, Path(path), (x, top, cell_w, img_h), alt)
        box = slide.shapes.add_textbox(x, top + img_h + Inches(0.04), cell_w,
                                       cap_h)
        tf = box.text_frame
        tf.word_wrap = True
        write(tf, [(caption, 10, SECOND, 0)])

    if note:
        box = slide.shapes.add_textbox(left, top + height - note_h, width,
                                       note_h)
        box.text_frame.word_wrap = True
        write(box.text_frame, [(note, 10, MUTED, 0)])
    return slide


def add_bullets_photo_slide(prs, title, bullets, photo, caption, alt):
    """Bullets on the left, one photograph on the right.

    The text goes into the layout's own body placeholder, narrowed, rather than
    into a new textbox — a textbox inherits neither the master's bullet glyphs
    nor its sans face, and the difference shows next to the other bullet slides.
    """
    layout = find_layout(prs, "title and content", "content", fallback=1)
    slide = prs.slides.add_slide(layout)
    style_title(slide, title)
    left, top, width, height = content_area(slide, prs)

    text_w = int(width * 0.58)
    img_x = left + text_w + Inches(0.25)
    img_w = width - text_w - Inches(0.25)

    bodies = body_placeholders(slide)
    if bodies:
        box = bodies[0]
        # Set all four, so python-pptx writes a complete xfrm rather than
        # inventing an offset of (0, 0) alongside a size.
        box.left, box.top, box.width, box.height = left, top, text_w, height
    else:
        box = slide.shapes.add_textbox(left, top, text_w, height)

    weight = sum(1 + len(t) // 55 for _, t in bullets)
    size = 13 if weight <= 9 else 12
    write(box.text_frame,
          [(t, size if lv == 0 else size - 1.5,
            INK if lv == 0 else MUTED, lv) for lv, t in bullets])
    clear_empty_placeholders(slide)

    cap_h = Inches(0.6)
    place_image(slide, prs, Path(photo), (img_x, top, img_w, height - cap_h),
                alt)
    cbox = slide.shapes.add_textbox(img_x, top + height - cap_h, img_w, cap_h)
    cbox.text_frame.word_wrap = True
    write(cbox.text_frame, [(caption, 9.5, SECOND, 0)])
    return slide


def add_groups_slide(prs, title, note=None):
    """One row per group of recommendations, with the numbers it contains."""
    groups = recommendation_groups()
    rows = [[name, ", ".join(nums)] for name, nums in groups]
    total = sum(len(nums) for _, nums in groups)
    note = (note or "") + (" %d recommendations in %d groups."
                           % (total, len(groups)))
    return add_table_slide(prs, title, ["Group", "Recommendations"], rows,
                           notes=[note.strip()])


def add_footers(prs, skip_first=True):
    """Place the footer line and slide number where the template puts them.

    Where the layout defines no footer position, nothing is added: the master
    already carries its own bottom chrome, and a textbox guessed at the bottom
    left would land on top of the logo.
    """
    total = len(prs.slides._sldIdLst)
    for i, slide in enumerate(prs.slides, start=1):
        if skip_first and i == 1:
            continue

        footer_rect, number_rect = layout_chrome(slide)

        # The template stamps the American Red Cross logo alone at the bottom
        # left of every content slide, from the slide master. Equal billing means
        # that arrangement has to go: the master mark is covered, the footer text
        # is dropped, and all four logos run across the band at a common height.
        # The slide number keeps its place at the right. Approved 2026-08-31.
        if logo_files():
            cover_master_logo(slide, prs)
            place_logo_row(slide, prs, Inches(0.34),
                           prs.slide_height - Inches(0.62), Inches(6.5),
                           Inches(0.40))
        elif footer_rect is not None:
            box = slide.shapes.add_textbox(*footer_rect)
            write(box.text_frame, [(FOOTER, 8, MUTED, 0)])
            box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        if number_rect is not None:
            box = slide.shapes.add_textbox(*number_rect)
            write(box.text_frame, [("%d / %d" % (i, total), 8, MUTED, 0)])
            box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def build(template=None, out_path="pdf/REPLICATION_RECOMMENDATIONS_BRIEFING.pptx"):
    prs = Presentation(template) if template else Presentation()

    # The templates ship with example slides. Start from an empty deck, dropping
    # the relationship as well as the list entry so no orphan part is left.
    for sld_id in list(prs.slides._sldIdLst):
        rid = sld_id.rId
        prs.slides._sldIdLst.remove(sld_id)
        prs.part.drop_rel(rid)

    for kind, payload in DECK:
        if kind == "title":
            add_title_slide(prs)
        elif kind == "section":
            add_section_slide(prs, payload["title"], payload.get("kicker", ""))
        elif kind == "bullets":
            add_bullets_slide(prs, payload["title"], payload["bullets"])
        elif kind == "table":
            add_table_slide(prs, payload["title"], payload["columns"],
                            payload["rows"], payload.get("notes"))
        elif kind == "groups":
            add_groups_slide(prs, payload["title"], payload.get("note"))
        elif kind == "figure":
            add_figure_slide(prs, payload["title"], payload["image"],
                             payload["alt"], payload.get("note"))
        elif kind == "photos":
            add_photos_slide(prs, payload["title"], payload["photos"],
                             payload.get("note"))
        elif kind == "bullets_photo":
            add_bullets_photo_slide(prs, payload["title"], payload["bullets"],
                                    payload["photo"], payload["caption"],
                                    payload["alt"])
        else:
            raise ValueError("unknown slide kind: %s" % kind)

    add_footers(prs)

    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out, len(prs.slides._sldIdLst)


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--template", default=DEFAULT_TEMPLATE,
                    help="Base .pptx or .potx to inherit theme and layouts from. "
                         "Defaults to the American Red Cross Classic template.")
    ap.add_argument("--no-template", action="store_true",
                    help="Build on python-pptx's neutral template instead.")
    ap.add_argument("-o", "--output",
                    default="pdf/REPLICATION_RECOMMENDATIONS_BRIEFING.pptx")
    args = ap.parse_args()

    template = None if args.no_template else args.template
    if template and not Path(template).is_file():
        sys.exit("Template not found: %s" % template)

    out, n = build(template, args.output)
    print("Wrote %s (%d slides) — %s" % (
        out, n, Path(template).name if template else "neutral template"))


if __name__ == "__main__":
    main()
